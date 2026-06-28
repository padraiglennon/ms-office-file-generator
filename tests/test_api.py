"""Route tests for the JSON API (ADR-007)."""

from __future__ import annotations

import io
import json

import pytest
from pptx import Presentation
from pptx.util import Inches

# Skip the whole module if the optional [web] extra is not installed.
pytest.importorskip("fastapi")

from fastapi.testclient import TestClient  # noqa: E402

from common_file_generator.web.app import create_app  # noqa: E402
from common_file_generator.web.caps import Caps  # noqa: E402

_OFFICE = "application/vnd.openxmlformats-officedocument"


@pytest.fixture
def client():
    return TestClient(create_app(max_upload_mb=25))


@pytest.fixture
def template_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.paragraphs[0].add_run().text = "{{title}}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def config_bytes() -> bytes:
    return json.dumps({"text": {"title": "Hello API"}}).encode()


def test_deck_streams_pptx_bytes(client) -> None:
    resp = client.post(
        "/api/generate/deck",
        json={"complexity": "simple", "slides": 3, "seed": 1, "background": "theme"},
    )
    assert resp.status_code == 200
    assert resp.headers["content-type"] == f"{_OFFICE}.presentationml.presentation"
    assert "attachment" in resp.headers["content-disposition"]
    assert "deck.pptx" in resp.headers["content-disposition"]
    # Opens as a real presentation.
    prs = Presentation(io.BytesIO(resp.content))
    assert len(prs.slides) == 3


def test_doc_streams_docx_bytes(client) -> None:
    from docx import Document

    resp = client.post(
        "/api/generate/doc",
        json={"complexity": "complex", "sections": 5, "seed": 2},
    )
    assert resp.status_code == 200
    assert resp.headers["content-type"] == f"{_OFFICE}.wordprocessingml.document"
    doc = Document(io.BytesIO(resp.content))
    assert len(doc.paragraphs) > 0


def test_sheet_streams_xlsx_bytes(client) -> None:
    from openpyxl import load_workbook

    resp = client.post(
        "/api/generate/sheet",
        json={"complexity": "maximum", "sheets": 3, "seed": 2},
    )
    assert resp.status_code == 200
    assert resp.headers["content-type"] == f"{_OFFICE}.spreadsheetml.sheet"
    wb = load_workbook(io.BytesIO(resp.content))
    assert len(wb.sheetnames) == 3


def test_pdf_streams_pdf_bytes(client) -> None:
    resp = client.post(
        "/api/generate/pdf",
        json={"complexity": "complex", "sections": 5, "seed": 2},
    )
    assert resp.status_code == 200
    assert resp.headers["content-type"] == "application/pdf"
    assert resp.content[:4] == b"%PDF"


def test_markdown_streams_md_bytes(client) -> None:
    resp = client.post(
        "/api/generate/markdown",
        json={"complexity": "complex", "sections": 5, "seed": 2},
    )
    assert resp.status_code == 200
    assert resp.headers["content-type"].startswith("text/markdown")
    assert resp.content.startswith(b"# ")


def test_minimal_body_uses_defaults(client) -> None:
    # All fields optional: an empty body generates with core defaults.
    resp = client.post("/api/generate/doc", json={})
    assert resp.status_code == 200
    assert len(resp.content) > 0


def test_doc_theme_streams_themed_docx(client) -> None:
    from docx import Document

    from common_file_generator.generators.docx_theme import SLATE

    resp = client.post(
        "/api/generate/doc",
        json={"complexity": "standard", "sections": 3, "seed": 2, "theme": "slate"},
    )
    assert resp.status_code == 200
    doc = Document(io.BytesIO(resp.content))
    assert str(doc.styles["Heading 1"].font.color.rgb) == str(SLATE.heading_color)


def test_unknown_theme_returns_422(client) -> None:
    # A bad theme is a request-shape error (rejected by the schema), not a 500.
    resp = client.post("/api/generate/doc", json={"theme": "neon"})
    assert resp.status_code == 422


def test_doc_request_theme_literal_tracks_themes() -> None:
    # Drift guard: the schema's allowed theme set must match the generator's.
    from typing import get_args

    from common_file_generator.generators.docx_theme import THEMES
    from common_file_generator.web.schemas import DocRequest

    allowed = set(get_args(DocRequest.model_fields["theme"].annotation))
    assert allowed == set(THEMES)


def test_invalid_complexity_returns_422(client) -> None:
    resp = client.post("/api/generate/doc", json={"complexity": "nonsense"})
    assert resp.status_code == 422
    body = resp.json()
    assert "detail" in body


def test_below_lower_bound_returns_400(client) -> None:
    # Value bounds (sections >= 1, cols >= 3, ...) are left to the core, so a
    # too-small count surfaces as a 400 with the core's message, matching the CLI.
    resp = client.post("/api/generate/doc", json={"sections": 0})
    assert resp.status_code == 400
    assert "at least 1" in resp.json()["detail"]


def test_cols_below_minimum_returns_400(client) -> None:
    resp = client.post("/api/generate/sheet", json={"cols": 2})
    assert resp.status_code == 400
    assert "at least 3" in resp.json()["detail"]


def test_unknown_background_returns_400(client) -> None:
    # background is a free string validated by the core, so a bad value is a
    # generation-time error surfaced as 400 with a JSON detail.
    resp = client.post("/api/generate/deck", json={"slides": 2, "background": "nope"})
    assert resp.status_code == 400
    assert "detail" in resp.json()


def test_fill_streams_file_and_report_header(
    client, template_bytes, config_bytes
) -> None:
    resp = client.post(
        "/api/generate/fill",
        files={
            "template": ("t.pptx", template_bytes, "application/octet-stream"),
            "config": ("c.json", config_bytes, "application/json"),
        },
    )
    assert resp.status_code == 200
    assert "filled.pptx" in resp.headers["content-disposition"]
    assert resp.headers.get("x-injection-report")  # report folded into a header
    # The filled file is a real presentation.
    Presentation(io.BytesIO(resp.content))


def test_fill_include_report_returns_full_json(
    client, template_bytes, config_bytes
) -> None:
    import base64

    resp = client.post(
        "/api/generate/fill?include_report=true",
        files={
            "template": ("t.pptx", template_bytes, "application/octet-stream"),
            "config": ("c.json", config_bytes, "application/json"),
        },
    )
    assert resp.status_code == 200
    body = resp.json()
    assert body["filename"] == "filled.pptx"
    # The full multi-line report survives intact (newlines preserved).
    assert "\n" in body["report"]
    assert "successfully" in body["report"] or "Injection report" in body["report"]
    # The file round-trips from base64 and opens.
    Presentation(io.BytesIO(base64.b64decode(body["file_base64"])))


def test_fill_rejects_wrong_template_type(client, config_bytes) -> None:
    resp = client.post(
        "/api/generate/fill",
        files={
            "template": ("t.txt", b"nope", "text/plain"),
            "config": ("c.json", config_bytes, "application/json"),
        },
    )
    assert resp.status_code == 400
    assert "not an accepted file type" in resp.json()["detail"]


def test_fill_rejects_oversize_upload(template_bytes, config_bytes) -> None:
    client = TestClient(create_app(max_upload_mb=1))
    big = b"x" * (2 * 1024 * 1024)
    resp = client.post(
        "/api/generate/fill",
        files={
            "template": ("t.pptx", big, "application/octet-stream"),
            "config": ("c.json", config_bytes, "application/json"),
        },
    )
    assert resp.status_code == 400
    assert "larger than" in resp.json()["detail"]


def test_openapi_documents_the_api(client) -> None:
    spec = client.get("/openapi.json").json()
    assert "/api/generate/deck" in spec["paths"]
    assert "/api/generate/fill" in spec["paths"]


def test_api_leaves_no_files_behind(client) -> None:
    # The streaming path uses a throwaway temp dir; nothing accumulates in the
    # UI's artifact registry.
    app = create_app()
    c = TestClient(app)
    c.post("/api/generate/markdown", json={"sections": 2})
    assert app.state.artifacts == {}


# --- Resource caps and runtime guards (ADR-010) ---


def _capped_client(**caps_kwargs) -> TestClient:
    return TestClient(create_app(caps=Caps(**caps_kwargs)))


def test_over_per_field_cap_returns_422() -> None:
    c = _capped_client(max_slides=5)
    resp = c.post("/api/generate/deck", json={"slides": 6})
    assert resp.status_code == 422
    detail = resp.json()["detail"]
    assert "slides" in detail and "maximum of 5" in detail


def test_within_per_field_cap_still_generates() -> None:
    c = _capped_client(max_slides=5)
    resp = c.post("/api/generate/deck", json={"slides": 3, "complexity": "minimal"})
    assert resp.status_code == 200


def test_composite_cost_over_budget_returns_422() -> None:
    # Each field is under its per-field cap, but sheets*tables*rows*cols exceeds
    # the composite budget.
    c = _capped_client(max_sheets=100, max_rows=5000, max_cols=100, max_cost=1000)
    resp = c.post(
        "/api/generate/sheet",
        json={"sheets": 5, "rows": 50, "cols": 10, "complexity": "minimal"},
    )
    assert resp.status_code == 422
    assert "cost" in resp.json()["detail"]


def test_composite_cost_uses_complexity_defaults() -> None:
    # No rows/cols supplied: the cost uses the per-complexity defaults, so a tiny
    # budget still rejects.
    c = _capped_client(max_cost=10)
    resp = c.post("/api/generate/doc", json={"sections": 20, "complexity": "maximum"})
    assert resp.status_code == 422


def test_generation_timeout_returns_503() -> None:
    # A zero-second budget trips the timeout guard before any real work finishes.
    c = _capped_client(gen_timeout_s=0)
    resp = c.post("/api/generate/doc", json={"sections": 1})
    assert resp.status_code == 503
    assert "time limit" in resp.json()["detail"]


def test_output_too_large_returns_400() -> None:
    # A 0 MB output cap rejects any non-empty file after generation.
    c = _capped_client(max_output_mb=0)
    resp = c.post("/api/generate/markdown", json={"sections": 1})
    assert resp.status_code == 400
    assert "exceeds" in resp.json()["detail"]


def test_fill_respects_timeout_guard(template_bytes, config_bytes) -> None:
    c = TestClient(create_app(caps=Caps(gen_timeout_s=0)))
    resp = c.post(
        "/api/generate/fill",
        files={
            "template": ("t.pptx", template_bytes, "application/octet-stream"),
            "config": ("c.json", config_bytes, "application/json"),
        },
    )
    assert resp.status_code == 503


# --- Global concurrency cap (ADR-013) ---


def _single_slot_app():
    # One slot, zero wait: with the slot held, the next request must 503 at once.
    return create_app(caps=Caps(max_concurrent=1, acquire_timeout_s=0))


def test_concurrency_full_returns_503() -> None:
    app = _single_slot_app()
    c = TestClient(app)
    # Occupy the only generation slot, then a request finds none free.
    assert app.state.limiter.acquire(timeout=0)
    try:
        resp = c.post("/api/generate/markdown", json={"sections": 1})
    finally:
        app.state.limiter.release()
    assert resp.status_code == 503
    assert "busy" in resp.json()["detail"]


def test_slot_released_after_success() -> None:
    # A normal request returns its slot, so the next one still succeeds.
    app = _single_slot_app()
    c = TestClient(app)
    assert c.post("/api/generate/markdown", json={"sections": 1}).status_code == 200
    assert c.post("/api/generate/markdown", json={"sections": 1}).status_code == 200


def test_slot_released_after_error() -> None:
    # An over-cap request (422) must not consume a slot; the limiter is acquired
    # only inside run_guarded, after the per-field check, but a failed generation
    # path must still release. Force a generation error and confirm capacity.
    app = create_app(caps=Caps(max_concurrent=1, acquire_timeout_s=0, max_output_mb=0))
    c = TestClient(app)
    # max_output_mb=0 makes generation raise OutputTooLarge (400) after building.
    assert c.post("/api/generate/markdown", json={"sections": 1}).status_code == 400
    # The slot was released despite the error: a second attempt reaches the same
    # guard rather than 503-ing on a leaked slot.
    assert c.post("/api/generate/markdown", json={"sections": 1}).status_code == 400


def test_fill_respects_concurrency_cap(template_bytes, config_bytes) -> None:
    app = _single_slot_app()
    c = TestClient(app)
    assert app.state.limiter.acquire(timeout=0)
    try:
        resp = c.post(
            "/api/generate/fill",
            files={
                "template": ("t.pptx", template_bytes, "application/octet-stream"),
                "config": ("c.json", config_bytes, "application/json"),
            },
        )
    finally:
        app.state.limiter.release()
    assert resp.status_code == 503
