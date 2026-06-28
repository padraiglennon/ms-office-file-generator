"""Route tests for the FastAPI + HTMX UI (ADR-002)."""

from __future__ import annotations

import io
import json
import re

import pytest
from pptx import Presentation
from pptx.util import Inches

# Skip the whole module if the optional [web] extra is not installed.
pytest.importorskip("fastapi")

from fastapi.testclient import TestClient  # noqa: E402

from common_file_generator.web.app import create_app  # noqa: E402
from common_file_generator.web.caps import Caps  # noqa: E402
from common_file_generator.web.forms import deck_fields  # noqa: E402


@pytest.fixture
def client():
    return TestClient(create_app(max_upload_mb=25))


@pytest.fixture
def template_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.paragraphs[0].add_run().text = "{{title}}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def config_bytes() -> bytes:
    return json.dumps({"text": {"title": "Hello UI"}}).encode()


def _token(html: str) -> str:
    match = re.search(r"/download/([\w-]+)", html)
    assert match, "no download link in response"
    return match.group(1)


def test_health_returns_ok(client) -> None:
    resp = client.get("/health")
    assert resp.status_code == 200
    assert resp.json() == {"status": "ok"}


def test_index_renders_introspected_fields(client) -> None:
    resp = client.get("/")
    assert resp.status_code == 200
    # complexity choices come from the Complexity enum via introspection
    assert "complexity" in resp.text
    assert "maximum" in resp.text
    assert "/static/htmx.min.js" in resp.text


def test_deck_fields_track_the_core() -> None:
    names = {f.name for f in deck_fields()}
    assert {"complexity", "slides", "seed", "background"} <= names
    complexity = next(f for f in deck_fields() if f.name == "complexity")
    assert "maximum" in {c.value for c in complexity.choices}
    background = next(f for f in deck_fields() if f.name == "background")
    assert "custom" in {c.value for c in background.choices}  # UI-only option


def test_every_deck_field_has_help_text() -> None:
    assert all(f.help.strip() for f in deck_fields())


def test_fields_render_tooltips_and_required_markers(client) -> None:
    from common_file_generator.web.forms import (
        doc_fields,
        md_fields,
        pdf_fields,
        sheet_fields,
    )

    html = client.get("/").text
    # Every field across the deck, document, spreadsheet, PDF, Markdown, and fill
    # forms carries a hover tooltip marker; help text is in its data-tip attribute.
    expected = (
        len(deck_fields())
        + len(doc_fields())
        + len(sheet_fields())
        + len(pdf_fields())
        + len(md_fields())
        + 2  # fill uploads
    )
    assert html.count('class="tip"') == expected
    assert "How busy each slide is" in html  # plain-language help (data-tip)
    assert 'class="help"' not in html  # no always-visible help lines
    # Required uploads are clearly marked.
    assert 'class="req"' in html
    assert "required" in html.lower()


def test_generate_deck_returns_downloadable_file(client) -> None:
    resp = client.post(
        "/generate/deck",
        data={
            "complexity": "simple",
            "slides": "3",
            "seed": "1",
            "background": "theme",
            "background_color": "",
            "video_url": "",
        },
    )
    assert resp.status_code == 200
    token = _token(resp.text)
    download = client.get(f"/download/{token}")
    assert download.status_code == 200
    assert len(download.content) > 0


def test_result_shows_timestamp_and_flash(client) -> None:
    resp = client.post(
        "/generate/deck",
        data={
            "complexity": "minimal",
            "slides": "2",
            "seed": "1",
            "background": "none",
        },
    )
    assert re.search(r"Generated at \d\d:\d\d:\d\d", resp.text)
    assert "flash" in resp.text  # animates on each swap so repeats are obvious


def test_fill_returns_file_and_report(client, template_bytes, config_bytes) -> None:
    resp = client.post(
        "/generate/fill",
        files={
            "template": ("t.pptx", template_bytes, "application/octet-stream"),
            "config": ("c.json", config_bytes, "application/json"),
        },
    )
    assert resp.status_code == 200
    assert "/download/" in resp.text
    assert "successfully" in resp.text or "Report" in resp.text


def test_fill_rejects_wrong_template_type(client, config_bytes) -> None:
    resp = client.post(
        "/generate/fill",
        files={
            "template": ("t.txt", b"nope", "text/plain"),
            "config": ("c.json", config_bytes, "application/json"),
        },
    )
    assert resp.status_code == 400
    assert "not an accepted file type" in resp.text


def test_fill_rejects_oversize_upload(template_bytes, config_bytes) -> None:
    client = TestClient(create_app(max_upload_mb=1))
    big = b"x" * (2 * 1024 * 1024)
    resp = client.post(
        "/generate/fill",
        files={
            "template": ("t.pptx", big, "application/octet-stream"),
            "config": ("c.json", config_bytes, "application/json"),
        },
    )
    assert resp.status_code == 400
    assert "larger than" in resp.text


def test_index_has_tabs(client) -> None:
    html = client.get("/").text
    assert 'role="tablist"' in html
    assert 'id="tab-deck"' in html and 'id="tab-fill"' in html
    assert 'id="panel-deck"' in html and 'id="panel-fill"' in html
    assert "showPanel(" in html  # the toggle script


def test_index_has_header_nav_links(client) -> None:
    html = client.get("/").text
    assert 'class="header-nav"' in html
    assert 'href="/docs"' in html  # FastAPI Swagger UI
    assert "padraiglennon.github.io/common-file-generator" in html  # Pages docs site


def test_index_has_document_tab(client) -> None:
    html = client.get("/").text
    assert 'id="tab-doc"' in html and 'id="panel-doc"' in html
    assert "Generate a document" in html
    assert 'hx-post="/generate/doc"' in html


def test_generate_document_returns_downloadable_docx(client) -> None:
    from docx import Document

    resp = client.post(
        "/generate/doc",
        data={"complexity": "complex", "sections": "5", "seed": "2"},
    )
    assert resp.status_code == 200
    download = client.get(f"/download/{_token(resp.text)}")
    assert download.status_code == 200
    doc = Document(io.BytesIO(download.content))
    assert len(doc.paragraphs) > 0


def test_generate_document_applies_selected_theme(client) -> None:
    from docx import Document

    from common_file_generator.generators.docx_theme import SLATE

    resp = client.post(
        "/generate/doc",
        data={"complexity": "standard", "sections": "3", "seed": "2", "theme": "slate"},
    )
    assert resp.status_code == 200
    download = client.get(f"/download/{_token(resp.text)}")
    doc = Document(io.BytesIO(download.content))
    assert str(doc.styles["Heading 1"].font.color.rgb) == str(SLATE.heading_color)


def test_generate_document_invalid_sections_errors(client) -> None:
    resp = client.post(
        "/generate/doc",
        data={"complexity": "standard", "sections": "0", "seed": "0"},
    )
    assert resp.status_code == 400
    assert "Could not generate" in resp.text


def test_ui_generation_timeout_returns_503() -> None:
    # The runtime guards (ADR-010) apply to the UI path too, via the shared
    # service. A zero-second budget trips the timeout.
    c = TestClient(create_app(caps=Caps(gen_timeout_s=0)))
    resp = c.post(
        "/generate/doc",
        data={"complexity": "standard", "sections": "1", "seed": "0"},
    )
    assert resp.status_code == 503


def test_ui_output_too_large_errors() -> None:
    c = TestClient(create_app(caps=Caps(max_output_mb=0)))
    resp = c.post(
        "/generate/md",
        data={"complexity": "standard", "sections": "1", "seed": "0"},
    )
    assert resp.status_code == 400


def test_ui_concurrency_full_returns_503() -> None:
    # The ADR-013 cap is shared with the API via the same service seam, so the UI
    # path 503s when no slot is free too.
    app = create_app(caps=Caps(max_concurrent=1, acquire_timeout_s=0))
    c = TestClient(app)
    assert app.state.limiter.acquire(timeout=0)
    try:
        resp = c.post(
            "/generate/md",
            data={"complexity": "standard", "sections": "1", "seed": "0"},
        )
    finally:
        app.state.limiter.release()
    assert resp.status_code == 503


def test_doc_fields_track_the_core() -> None:
    from common_file_generator.web.forms import doc_fields

    names = {f.name for f in doc_fields()}
    assert {"complexity", "sections", "seed", "theme"} == names


def test_doc_theme_field_offers_the_three_themes() -> None:
    from common_file_generator.generators.docx_theme import THEMES
    from common_file_generator.web.forms import doc_fields

    theme_field = next(f for f in doc_fields() if f.name == "theme")
    assert theme_field.kind == "select"
    assert {c.value for c in theme_field.choices} == set(THEMES)


def test_pdf_and_md_forms_have_no_theme_field() -> None:
    from common_file_generator.web.forms import md_fields, pdf_fields

    assert "theme" not in {f.name for f in pdf_fields()}
    assert "theme" not in {f.name for f in md_fields()}


def test_index_has_spreadsheet_tab(client) -> None:
    html = client.get("/").text
    assert 'id="tab-sheet"' in html and 'id="panel-sheet"' in html
    assert "Generate a spreadsheet" in html
    assert 'hx-post="/generate/sheet"' in html


def test_generate_spreadsheet_returns_downloadable_xlsx(client) -> None:
    from openpyxl import load_workbook

    resp = client.post(
        "/generate/sheet",
        data={"complexity": "maximum", "sheets": "3", "seed": "2"},
    )
    assert resp.status_code == 200
    download = client.get(f"/download/{_token(resp.text)}")
    assert download.status_code == 200
    wb = load_workbook(io.BytesIO(download.content))
    assert len(wb.sheetnames) == 3


def test_generate_spreadsheet_invalid_sheets_errors(client) -> None:
    resp = client.post(
        "/generate/sheet",
        data={"complexity": "standard", "sheets": "0", "seed": "0"},
    )
    assert resp.status_code == 400
    assert "Could not generate" in resp.text


def test_sheet_fields_track_the_core() -> None:
    from common_file_generator.web.forms import sheet_fields

    names = {f.name for f in sheet_fields()}
    assert {"complexity", "sheets", "seed"} == names


def test_index_has_pdf_tab(client) -> None:
    html = client.get("/").text
    assert 'id="tab-pdf"' in html and 'id="panel-pdf"' in html
    assert "Generate a PDF" in html
    assert 'hx-post="/generate/pdf"' in html


def test_index_has_markdown_tab(client) -> None:
    html = client.get("/").text
    assert 'id="tab-md"' in html and 'id="panel-md"' in html
    assert "Generate Markdown" in html
    assert 'hx-post="/generate/md"' in html


def test_generate_pdf_returns_downloadable_file(client) -> None:
    resp = client.post(
        "/generate/pdf",
        data={"complexity": "complex", "sections": "5", "seed": "2"},
    )
    assert resp.status_code == 200
    download = client.get(f"/download/{_token(resp.text)}")
    assert download.status_code == 200
    assert download.content[:4] == b"%PDF"


def test_generate_pdf_invalid_sections_errors(client) -> None:
    resp = client.post(
        "/generate/pdf",
        data={"complexity": "standard", "sections": "0", "seed": "0"},
    )
    assert resp.status_code == 400
    assert "Could not generate" in resp.text


def test_generate_markdown_returns_downloadable_md(client) -> None:
    resp = client.post(
        "/generate/md",
        data={"complexity": "complex", "sections": "5", "seed": "2"},
    )
    assert resp.status_code == 200
    download = client.get(f"/download/{_token(resp.text)}")
    assert download.status_code == 200
    assert download.content.startswith(b"# ")


def test_generate_markdown_invalid_sections_errors(client) -> None:
    resp = client.post(
        "/generate/md",
        data={"complexity": "standard", "sections": "0", "seed": "0"},
    )
    assert resp.status_code == 400
    assert "Could not generate" in resp.text


def test_pdf_fields_track_the_core() -> None:
    from common_file_generator.web.forms import pdf_fields

    names = {f.name for f in pdf_fields()}
    assert {"complexity", "sections", "seed"} == names


def test_md_fields_track_the_core() -> None:
    from common_file_generator.web.forms import md_fields

    names = {f.name for f in md_fields()}
    assert {"complexity", "sections", "seed"} == names


def test_background_control_is_merged(client) -> None:
    html = client.get("/").text
    # one Slide background control with friendly labels + a Custom colour option
    assert "Plain white" in html
    assert "Custom colour..." in html
    assert 'type="color"' in html  # the revealed picker
    assert "Custom background colour (optional)" not in html  # no separate field


def test_custom_background_applies_chosen_colour(client) -> None:
    from pptx import Presentation

    resp = client.post(
        "/generate/deck",
        data={
            "complexity": "simple",
            "slides": "3",
            "seed": "1",
            "background": "custom",
            "background_color": "#ff0000",
        },
    )
    assert resp.status_code == 200
    download = client.get(f"/download/{_token(resp.text)}")
    prs = Presentation(io.BytesIO(download.content))
    # Every slide, including the title slide, gets the chosen colour.
    tints = {str(slide.background.fill.fore_color.rgb) for slide in prs.slides}
    assert tints == {"FF0000"}


def test_deck_accepts_hash_prefixed_colour(client) -> None:
    resp = client.post(
        "/generate/deck",
        data={
            "complexity": "simple",
            "slides": "3",
            "seed": "1",
            "background": "custom",
            "background_color": "#ff0000",
        },
    )
    assert resp.status_code == 200
    assert "/download/" in resp.text


def test_deck_invalid_value_shows_error(client) -> None:
    resp = client.post(
        "/generate/deck",
        data={
            "complexity": "nonsense",
            "slides": "3",
            "seed": "0",
            "background": "none",
            "background_color": "",
            "video_url": "",
        },
    )
    assert resp.status_code == 400
    assert "Could not generate" in resp.text


def test_download_unknown_token_404(client) -> None:
    assert client.get("/download/nope").status_code == 404


def test_expired_artifacts_are_swept() -> None:
    # ttl_seconds=0 means any prior artifact is past its TTL on the next store.
    app = create_app(ttl_seconds=0)
    c = TestClient(app)
    first = c.post(
        "/generate/deck",
        data={
            "complexity": "minimal",
            "slides": "2",
            "seed": "1",
            "background": "none",
            "background_color": "",
            "video_url": "",
        },
    )
    old_token = _token(first.text)
    old_path = app.state.artifacts[old_token].path
    assert old_path.is_file()

    # A second generation triggers the sweep, removing the expired first file.
    c.post(
        "/generate/deck",
        data={
            "complexity": "minimal",
            "slides": "2",
            "seed": "2",
            "background": "none",
            "background_color": "",
            "video_url": "",
        },
    )
    assert old_token not in app.state.artifacts
    assert not old_path.exists()
    assert c.get(f"/download/{old_token}").status_code == 404


def test_server_defaults_to_localhost_18990() -> None:
    from common_file_generator.web.server import build_parser

    args = build_parser().parse_args([])
    assert args.host == "127.0.0.1"
    assert args.port == 18990
    assert args.max_upload_mb == 25


def test_server_env_overrides_max_upload(monkeypatch) -> None:
    monkeypatch.setenv("COMMON_FILE_GEN_MAX_UPLOAD_MB", "100")
    import importlib

    from common_file_generator.web import server

    importlib.reload(server)
    args = server.build_parser().parse_args([])
    assert args.max_upload_mb == 100
    monkeypatch.delenv("COMMON_FILE_GEN_MAX_UPLOAD_MB")
    importlib.reload(server)


def test_server_host_override() -> None:
    from common_file_generator.web.server import build_parser

    args = build_parser().parse_args(["--host", "0.0.0.0", "--port", "9000"])
    assert args.host == "0.0.0.0"
    assert args.port == 9000


def test_server_env_overrides_host_and_port(monkeypatch) -> None:
    monkeypatch.setenv("COMMON_FILE_GEN_HOST", "0.0.0.0")
    monkeypatch.setenv("COMMON_FILE_GEN_PORT", "28990")
    import importlib

    from common_file_generator.web import server

    importlib.reload(server)
    args = server.build_parser().parse_args([])
    assert args.host == "0.0.0.0"
    assert args.port == 28990
    # An explicit flag still wins over the environment.
    flagged = server.build_parser().parse_args(["--port", "9000"])
    assert flagged.port == 9000
    monkeypatch.delenv("COMMON_FILE_GEN_HOST")
    monkeypatch.delenv("COMMON_FILE_GEN_PORT")
    importlib.reload(server)


def test_server_ignores_invalid_env_port(monkeypatch) -> None:
    monkeypatch.setenv("COMMON_FILE_GEN_PORT", "not-a-number")
    import importlib

    from common_file_generator.web import server

    importlib.reload(server)
    args = server.build_parser().parse_args([])
    assert args.port == 18990  # falls back to the default
    monkeypatch.delenv("COMMON_FILE_GEN_PORT")
    importlib.reload(server)
